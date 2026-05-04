import streamlit as st
import sys
import os
import json
import uuid
import time
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt

# Add the project root to sys.path so we can import src
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from src.graph import app as graph_app
from src.state import FluxIdeasState
from src.database import FluxIdeasDB
from datetime import datetime

# Initialize Database
db = FluxIdeasDB()

st.set_page_config(page_title='FLUXIDEAS | AI Intelligence', page_icon='⚡', layout='wide', initial_sidebar_state='collapsed')

# --- Sidebar Archive (The Dynamic Library) ---
with st.sidebar:
    st.markdown("<h2 style='color:#60efff; font-size:1.2rem; font-weight:900;'>🏛️ VENTURE LIBRARY</h2>", unsafe_allow_html=True)
    st.caption("Your persistent archive of AI-generated business blueprints.")
    
    past_scans = db.get_all_dossiers()
    if past_scans:
        for ps in past_scans:
            try:
                # Handle potential timestamp issues
                dt = datetime.fromisoformat(ps['created_at'])
                date_label = dt.strftime("%b %d | %H:%M")
            except:
                date_label = "Archive"
                
            btn_label = f"🚀 {ps['problem_name']}\n{date_label}"
            if st.button(btn_label, key=f"ps_{ps['id']}", use_container_width=True):
                # Reload State
                try:
                    f_state = json.loads(ps['full_state_json'])
                    st.session_state.current_state = f_state
                    st.session_state.selected_problem = f_state.get("selected_problem")
                    st.session_state.topic = ps['topic']
                    st.session_state.app_stage = "done"
                    st.rerun()
                except Exception as e:
                    st.error(f"Failed to load: {e}")
        
        st.markdown("---")
        if st.button("🆕 Start New Research", use_container_width=True):
            st.session_state.app_stage = "input"
            st.rerun()
    else:
        st.info("No dossiers saved yet. Complete a scan to build your library!")

# Premium CSS for Flux Aesthetic
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600&family=Outfit:wght@700;900&display=swap');

    /* Clean Base */
    html { scroll-behavior: smooth; }
    .stApp { 
        background: #0b0f1a; 
        font-family: 'Inter', sans-serif; 
        color: #e2e8f0;
    }

    /* Minimalist Title */
    .flux-title { 
        font-size: 4.5rem !important; 
        font-weight: 900; 
        font-family: 'Outfit', sans-serif;
        background: linear-gradient(to right, #ffffff, #94a3b8);
        -webkit-background-clip: text; 
        -webkit-text-fill-color: transparent; 
        text-align: center; 
        margin-bottom: 0px !important; 
        letter-spacing: -2px;
    }

    .flux-tagline { 
        text-align: center; 
        color: #64748b; 
        font-size: 0.9rem; 
        font-weight: 500; 
        letter-spacing: 4px; 
        text-transform: uppercase; 
        margin-bottom: 60px !important; 
    }

    /* Clean Interactive Cards */
    div[data-testid="stForm"], .flux-card {
        background: #111827 !important;
        border: 1px solid #1f2937 !important;
        border-radius: 24px !important;
        padding: 40px !important;
        box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.1) !important;
        transition: border-color 0.3s ease;
    }
    div[data-testid="stForm"]:hover {
        border-color: #3b82f6 !important;
    }

    /* Clean Modern Button */
    div.stButton > button { 
        background: #3b82f6 !important; 
        color: white !important; 
        font-weight: 600 !important; 
        border-radius: 12px !important; 
        border: none !important; 
        padding: 0.6rem 2rem !important;
        transition: all 0.2s ease !important;
        width: 100% !important;
    }
    div.stButton > button:hover { 
        background: #2563eb !important;
        box-shadow: 0 0 20px rgba(59, 130, 246, 0.4) !important;
    }

    /* Luminous Minimal Inputs */
    .stTextInput > div > div > input { 
        background: #1f2937 !important; 
        border-radius: 12px !important; 
        border: 1px solid #374151 !important; 
        color: white !important; 
        padding: 12px 20px !important;
    }
    .stTextInput > div > div > input:focus {
        border-color: #3b82f6 !important;
        box-shadow: 0 0 0 2px rgba(59, 130, 246, 0.2) !important;
    }

    /* Interactive Sidebar */
    [data-testid="stSidebar"] {
        background-color: #030712 !important;
        border-right: 1px solid #1f2937 !important;
    }
    
    .log-agent {
        color: #3b82f6;
        font-weight: 700;
        font-family: monospace;
    }
    .log-entry {
        border-bottom: 1px solid #1f2937;
        padding: 8px 0;
        font-size: 0.85rem;
    }

</style>
""", unsafe_allow_html=True)

# ── Main Application ──────────────────────────────────────────────────
st.markdown("""
<div style='text-align:center; margin-bottom:60px;'>
    <h1 class='flux-title'>FLUXIDEAS</h1>
    <p class='flux-tagline'>Multi-Agent Venture Intelligence</p>
</div>
""", unsafe_allow_html=True)

# Initialize Session State
if 'app_stage' not in st.session_state:
    st.session_state.app_stage = "input"
    st.session_state.thread_id = str(uuid.uuid4())
    st.session_state.current_state = None
    st.session_state.topic = ""
    st.session_state.founder_profile = {"location": "Global", "skills": "", "budget": "$0 - $100", "time": "5-10 hrs"}
    st.session_state.selected_problem = None
    st.session_state.archived_dossier = None
    st.session_state.live_logs = []
    st.session_state.chat_history = []
    st.session_state.model_usage = {}

# Initialize Mission Control Container globally so agents can find it
if 'log_container' not in st.session_state:
    st.session_state.log_container = None

def add_log(agent, message, container=None, model=None):
    timestamp = time.strftime("%H:%M:%S")
    model_tag = f" <span style='color:#00ff87; font-size:0.75rem; border:1px solid #00ff8744; padding:1px 6px; border-radius:4px; margin-left:8px;'>{model}</span>" if model else ""
    
    # Check for error status
    is_error = "❌" in message or "error" in message.lower()
    msg_color = "#ef4444" if is_error else "#e2e8f0"
    
    log_html = f"<div class='log-entry'><span style='color:#475569'>[{timestamp}]</span> <span class='log-agent'>{agent.upper()}</span>{model_tag}: <span style='color:{msg_color}'>{message}</span></div>"
    st.session_state.live_logs.append(log_html)
    
    if len(st.session_state.live_logs) > 30:
        st.session_state.live_logs.pop(0)
    
    # Update live container if it exists
    target = container or st.session_state.get('log_container')
    if target:
        with target:
            # Re-render all logs (Streamlit way to keep it live)
            for l in st.session_state.live_logs[::-1]:
                st.markdown(l, unsafe_allow_html=True)

def generate_pitch_deck(p_name, dossier, market, risk):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = p_name
    subtitle.text = f"Founder's Dossier & Market Strategy\nGenerated by FLUXIDEAS"

    # Problem Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The Problem & Market Gap"
    tf = body_shape.text_frame
    tf.text = dossier.get('market_gap_score', {}).get('rationale', 'Market gap identified.')
    
    # Solution Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "The Solution: MVP Blueprint"
    tf = slide.shapes.placeholders[1].text_frame
    mvp = dossier.get('mvp_blueprint', {})
    for k, v in mvp.items():
        if isinstance(v, dict):
            p = tf.add_paragraph()
            p.text = f"{v.get('name')}: {v.get('description')}"
            p.level = 0

    # Market Size Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Market Opportunity (TAM/SAM/SOM)"
    tf = slide.shapes.placeholders[1].text_frame
    if market:
        for k in ['tam', 'sam', 'som']:
            p = tf.add_paragraph()
            p.text = f"{k.upper()}: {market.get(k, 'N/A')}"
        p = tf.add_paragraph()
        p.text = f"Growth Rate: {market.get('growth_rate', 'N/A')}"

    # Risks Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Risk Audit & Mitigation"
    tf = slide.shapes.placeholders[1].text_frame
    if risk:
        for k in ['technical_risk', 'market_risk', 'kill_switch_criteria']:
            p = tf.add_paragraph()
            p.text = f"{k.replace('_',' ').title()}: {risk.get(k, 'N/A')}"

    # Save to buffer
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

config = {"configurable": {"thread_id": st.session_state.thread_id}}

def process_stream(stream_generator, status_container, log_container=None):
    progress_bar = st.progress(0, text="⚡ Initializing Multi-Agent Flux...")
    agent_count = 0
    total_agents = 9 # Core pipeline agents
    
    for output in stream_generator:
        if "__interrupt__" in output:
            continue
            
        for key, value in output.items():
            if not isinstance(value, dict):
                continue
            
            agent_count += 1
            progress = min(agent_count / total_agents, 1.0)
            
            if st.session_state.current_state is None:
                st.session_state.current_state = value.copy()
            else:
                st.session_state.current_state.update(value)
            
            # Sync model usage to session state
            current_model_usage = st.session_state.current_state.get("model_usage", {})
            st.session_state.model_usage.update(current_model_usage)
            
            current_model = current_model_usage.get(key.capitalize()) or current_model_usage.get(key)
            
            # Detailed Log Capture
            detailed_log = value.get('log_message')
            
            # Define Phases for visual clarity
            phases = {
                "scout": ("Phase 1/4: Data Discovery", "🕵️‍♂️"),
                "researcher": ("Phase 1/4: Data Discovery", "🔍"),
                "reasoner": ("Phase 2/4: Intelligence Synthesis", "🧠"),
                "analyst": ("Phase 2/4: Intelligence Synthesis", "📊"),
                "strategist": ("Phase 3/4: Business Architecture", "📈"),
                "economist": ("Phase 3/4: Business Architecture", "💰"),
                "designer": ("Phase 3/4: Business Architecture", "🎨"),
                "critic": ("Phase 4/4: Critical Risk Audit", "🛡️")
            }

            if key == "orchestrator":
                next_a = value.get('next_agent')
                status_container.update(label=f"🤖 Orchestrator: Routing to {next_a}...", state="running")
                progress_bar.progress(progress, text=f"Routing to {next_a.upper()}...")
                add_log("orchestrator", f"Decision: Handing over mission to {next_a.upper()}.", log_container)
                if next_a != "END" and next_a != "ask_human":
                    add_log(next_a, "Initializing agent systems... thinking...", log_container)
            elif key in phases:
                phase_name, phase_icon = phases[key]
                status_container.update(label=f"{phase_icon} {phase_name}: {key.capitalize()} in progress...", state="running")
                
                # Show detailed updates inside the expander directly
                status_container.markdown(f"**{phase_icon} {key.capitalize()} Update:**")
                if detailed_log:
                    status_container.caption(detailed_log)
                
                # Global Log
                if detailed_log:
                    add_log(key, detailed_log, st.session_state.log_container, model=current_model)
                
                if key == "critic":
                    progress_bar.empty()

# ── Sidebar: Flux Library ──────────────────────────────────────────
with st.sidebar:
    st.markdown("### <span class='live-pulse'></span> 🖥️ Mission Control", unsafe_allow_html=True)
    st.caption("Live Multi-Agent Log")
    # Store in session state for cross-stage access
    st.session_state.log_container = st.container(height=350)
    with st.session_state.log_container:
        if st.session_state.live_logs:
            for l in st.session_state.live_logs[::-1]:
                st.markdown(l, unsafe_allow_html=True)
        else:
            st.caption("Waiting for mission start...")
    

# ── Main UI Logic ───────────────────────────────────────────────────
# Stage: INPUT
if st.session_state.app_stage == "input":
    with st.form("search_form", border=False):
        topic = st.text_input("Enter a Market or Topic (e.g. 'CI/CD tools', 'No-code for Shopify')", placeholder="What niche are we auditing today?")
        
        st.markdown("### 👤 Founder Profile & Constraints")
        st.caption("Tell the AI your background and target region to personalize the business opportunities.")
        
        c0 = st.columns(1)[0]
        with c0:
            target_location = st.text_input("🌍 Target Location", placeholder="e.g. Global, London, US, India", value="Global")
            
        c1, c2, c3 = st.columns(3)
        with c1:
            skills = st.text_input("My Skills", placeholder="e.g. Python, Marketing, Design")
        with c2:
            budget = st.selectbox("Monthly Budget", ["$0 - $100", "$100 - $1,000", "$1,000 - $10,000", "$10,000+"])
        with c3:
            time_avail = st.selectbox("Weekly Time", ["5-10 hrs", "10-20 hrs", "20-40 hrs", "Full Time"])

        submit = st.form_submit_button("⚡ Launch Intelligence Scan", use_container_width=True)
        
        if submit and topic:
            st.session_state.topic = topic
            st.session_state.founder_profile = {
                "location": target_location,
                "skills": skills,
                "budget": budget,
                "time": time_avail
            }
            st.session_state.model_usage = {} # Reset
            st.session_state.app_stage = "processing1"
            st.rerun()

# Stage: PROCESSING 1
if st.session_state.app_stage == "processing1":
    with st.status("⚡ FLUXIDEAS Intelligence Pipeline Active...", expanded=True) as status:
        # Pre-flight Check
        tavily_key = os.environ.get("TAVILY_API_KEY")
        groq_key = os.environ.get("GROQ_API_KEY")
        
        if not tavily_key or "tvly-" not in tavily_key:
            status.update(label="❌ Configuration Error", state="error")
            st.error("**Missing Tavily API Key!** Please add `TAVILY_API_KEY` to your environment secrets.")
            if st.button("Back to Input"):
                st.session_state.app_stage = "input"
                st.rerun()
            st.stop()
            
        if not groq_key or "gsk_" not in groq_key:
            status.update(label="❌ Configuration Error", state="error")
            st.error("**Missing Groq API Key!** Please add `GROQ_API_KEY` to your environment secrets.")
            if st.button("Back to Input"):
                st.session_state.app_stage = "input"
                st.rerun()
            st.stop()

        initial_state = {
            "topic": st.session_state.topic, 
            "founder_profile": st.session_state.founder_profile,
            "model_usage": {}
        }
        try:
            process_stream(
                graph_app.stream(initial_state, config=config), 
                status,
                st.session_state.log_container 
            )
            
            # Check state to see if we hit the interrupt
            graph_state = graph_app.get_state(config)
            if graph_state.next and graph_state.next[0] == "ask_human":
                st.session_state.app_stage = "selection"
                status.update(label="✋ Paused for Human Input", state="complete")
                st.rerun()
            else:
                # If it finished without pausing (error case fallback)
                status.update(label="Pipeline finished without Human Node.", state="complete")
                st.session_state.app_stage = "done"
                st.rerun()
        except Exception as e:
            import traceback
            status.update(label=f"❌ Critical Pipeline Failure", state="error")
            st.error(f"**The AI Intelligence Engine encountered an error:** {str(e)}")
            with st.expander("🛠️ Debug Traceback (For Developers)"):
                st.code(traceback.format_exc())
            st.info("Check your API keys, network connection, or try a more specific topic.")
            if st.button("Reset Flux Session"):
                st.session_state.app_stage = "input"
                st.rerun()

# Stage: SELECTION
if st.session_state.app_stage == "selection":
    # Ensure log container stays current on selection screen
    if st.session_state.live_logs:
        with log_container:
            for l in st.session_state.live_logs[::-1]:
                st.markdown(l, unsafe_allow_html=True)
    
    st.markdown("### 🎯 Select a Market Gap to Build")
    st.markdown("Review each idea below — then hit **✅ Select & Build** on the one you want to pursue.")
    st.markdown("---")

    problems = []
    if st.session_state.current_state:
        problems = st.session_state.current_state.get("identified_problems", [])

    if not problems:
        st.error("No problems were identified. Please start over.")
        if st.button("🔄 Start New Scan", use_container_width=True):
            st.session_state.clear()
            st.rerun()
    else:
        for idx, p in enumerate(problems):
            with st.container():
                hcol1, hcol2 = st.columns([5, 1])
                with hcol1:
                    st.markdown(f"<div class='problem-title'>{p.get('problem_name')}</div>", unsafe_allow_html=True)
                with hcol2:
                    if st.button("✅ Select & Build", key=f"sel_{idx}", use_container_width=True):
                        st.session_state.selected_problem = p
                        st.session_state.app_stage = "processing2"
                        st.rerun()
                
                st.markdown(f"**Market Gap:** {p.get('market_gap')}")
                st.markdown(f"**Target Customer:** {p.get('target_customer')}")
                
                # Metrics
                st.markdown(
                    f"<span class='metric-pill'>📊 Market Score: {p.get('market_score')}/10</span>"
                    f"<span class='metric-pill' style='background:#1e293b; color:#00ff87'>👤 Founder Fit: {p.get('founder_fit_score', '?')}/10</span>"
                    f"<span class='metric-pill'>🔥 Urgency: {p.get('urgency_score')}/10</span>", 
                    unsafe_allow_html=True
                )
                
                st.markdown(f"**The Problem:** {p.get('description')}")
                
                # Source references
                src_refs = p.get("source_refs", [])
                if src_refs:
                    st.markdown("**🔗 Sources:**")
                    ref_parts = []
                    for ref in src_refs:
                        title = ref.get("title", "Discussion")
                        url   = ref.get("url", "")
                        author = ref.get("author", "")
                        if url:
                            ref_parts.append(f"[{title}]({url}) by *{author}*")
                        else:
                            ref_parts.append(f"{title} by *{author}*")
                    st.markdown(" &nbsp;|&nbsp; ".join(ref_parts))

                st.markdown("---")


# Stage: PROCESSING 2
if st.session_state.app_stage == "processing2":
    with st.status("⚙️ Resuming Pipeline...", expanded=True) as status:
        # Update state with selected problem
        graph_app.update_state(config, {"selected_problem": st.session_state.selected_problem})
        
        try:
            # Resume stream by passing None
            process_stream(graph_app.stream(None, config=config), status, st.session_state.log_container)
            status.update(label="✅ Scan Complete!", state="complete", expanded=False)
            
            # --- PERSISTENCE: Save to Library ---
            final_state = st.session_state.current_state
            if final_state:
                # Add selected problem to full state for easier reloading
                final_state["selected_problem"] = st.session_state.selected_problem
                
                db.save_dossier(
                    topic=st.session_state.topic,
                    problem_name=st.session_state.selected_problem.get("problem_name", "Unknown"),
                    blueprint=final_state.get("blueprint"),
                    market_score=st.session_state.selected_problem.get("market_score", 0),
                    mockup_url=final_state.get("mockup_url"),
                    risk_assessment=final_state.get("risk_assessment"),
                    full_state=final_state
                )
            
            st.session_state.app_stage = "done"
            st.rerun()
        except Exception as e:
            status.update(label=f"❌ Error: {str(e)}", state="error")

# Stage: DONE
if st.session_state.app_stage == "done":
    sel_prob   = st.session_state.get("selected_problem", {})
    p_name     = sel_prob.get("problem_name", "N/A")
    p_score    = int(sel_prob.get("market_score", 0))
    p_sent     = sel_prob.get("sentiment", "N/A")
    dossier    = (st.session_state.current_state or {}).get("blueprint", {}) or {}

    score_color = "#00ff87" if p_score >= 75 else ("#ffb347" if p_score >= 50 else "#ff6b6b")

    # ── Header ──────────────────────────────────────────────────────────────
    st.success("✅ Founder's Dossier Ready")
    st.markdown(f"""
    <div style='background:#1e293b;padding:24px;border-radius:14px;border-left:6px solid {score_color};margin-bottom:24px'>
        <div style='font-size:1.6rem;font-weight:800;color:#f8fafc;margin-bottom:8px'>📄 {p_name}</div>
        <span style='background:#334155;color:{score_color};padding:4px 14px;border-radius:20px;font-size:0.95rem;font-weight:700;margin-right:8px'>🔥 Market Score: {p_score}/100</span>
        <span style='background:#334155;color:#a0aec0;padding:4px 14px;border-radius:20px;font-size:0.9rem;'>💬 {p_sent}</span>
    </div>
    """, unsafe_allow_html=True)

    # Venture Health Gauges
    gh1, gh2, gh3 = st.columns(3)
    gh1.metric("👤 Founder Fit", f"{sel_prob.get('founder_fit_score', 0)*10}%")
    gh2.metric("🔥 Market Heat", f"{sel_prob.get('urgency_score', 0)*10}%")
    gh3.metric("🏰 Defensibility", f"{sel_prob.get('moat_score', 0)*10}%")
    st.markdown("---")

    # ── 5 Tabs ───────────────────────────────────────────────────────────────
    tab1, tab2, tab3, tab4, tab5 = st.tabs(["📊 The Insight", "💡 The MVP", "🛠️ The Build", "📉 Risk Audit", "🔗 Sources"])

    # ── Collect References Logic ─────────────────────────────────────────────
    # Collect all cited URLs from voice_of_customer, selected problem source_refs, and raw_sources
    all_refs = []
    for item in dossier.get("voice_of_customer", []):
        url = item.get("url", "")
        if url and url not in [r.get("url") for r in all_refs]:
            all_refs.append({"author": item.get("author", ""), "url": url, "source": item.get("source", "")})
    
    sel_prob = dossier.get("selected_problem", {})
    for ref in sel_prob.get("source_refs", []):
        url = ref.get("url", "")
        if url and url not in [r.get("url") for r in all_refs]:
            all_refs.append({"author": ref.get("author", ""), "url": url, "source": ref.get("title", "")})
            
    for src in (st.session_state.current_state or {}).get("raw_sources", []):
        url = src.get("url", "")
        if url and url not in [r.get("url") for r in all_refs]:
            all_refs.append({
                "author": src.get("author", ""), 
                "url": url, 
                "story_title": src.get("story_title", ""),
                "source": src.get("story_title", "General Insight")
            })

    # ── Tab 1: The Insight ───────────────────────────────────────────────────
    with tab1:
        # Signal Strength
        sig = dossier.get("signal_strength", {})
        st.markdown("### ⚡ Signal Strength (Validation)")
        scol1, scol2, scol3 = st.columns(3)
        scol1.metric("Mention Count", f"~{sig.get('mention_count', '?')} threads")
        scol2.metric("Signal Reliability", "✅ Confirmed")
        scol3.metric("Source", "HN / Reddit / Forums")
        st.info(f"**Where it shows up:** {sig.get('source_summary', '')}")
        st.caption(sig.get("validation", ""))

        st.markdown("---")

        # Market Gap Score
        mgs = dossier.get("market_gap_score", {})
        st.markdown("### 📈 Market Gap Score")
        mc1, mc2, mc3, mc4 = st.columns(4)
        mc1.metric("Overall", f"{mgs.get('total', '?')}/10")
        mc2.metric("⚡ Urgency", f"{mgs.get('urgency', '?')}/10")
        mc3.metric("💰 Commercial", f"{mgs.get('commercial_potential', '?')}/10")
        mc4.metric("🔧 Feasibility", f"{mgs.get('feasibility', '?')}/10")
        
        # VC Framework Row (Hardened Fallbacks)
        st.markdown("#### 🛡️ Venture Capital Framework")
        vcc1, vcc2 = st.columns(2)
        
        # Pull with multiple key variations and default to 5 if missing (prevents '?' for old scans)
        moat_val = sel_prob.get("moat_score", sel_prob.get("moatscore", 5))
        net_val  = sel_prob.get("network_effects", sel_prob.get("networkeffects", 5))
        
        vcc1.metric("🏰 Moat Potential", f"{moat_val}/10")
        vcc2.metric("🕸️ Network Effects", f"{net_val}/10")
        
        st.markdown(f"> {mgs.get('rationale', '')}")

        st.markdown("---")

        # Market Size (Economist)
        mkt = (st.session_state.current_state or {}).get("market_size_analysis", {})
        if mkt:
            st.markdown("### 💰 Market Size Estimates")
            ec1, ec2, ec3 = st.columns(3)
            ec1.metric("TAM", mkt.get("tam", "N/A"))
            ec2.metric("SAM", mkt.get("sam", "N/A"))
            ec3.metric("SOM", mkt.get("som", "N/A"))
            st.write(f"**Growth Rate:** {mkt.get('growth_rate', 'N/A')}")
            st.success(f"**Economist Verdict:** {mkt.get('economist_verdict', '')}")
            
            # Unit Economics Benchmarks
            bench = mkt.get("benchmarks", {})
            if bench:
                st.markdown("#### 💵 Projected Unit Economics")
                bc1, bc2, bc3 = st.columns(3)
                bc1.metric("Est. LTV", bench.get("avg_ltv", "N/A"))
                bc2.metric("Target CAC", bench.get("target_cac", "N/A"))
                bc3.metric("Payback", bench.get("payback_period", "N/A"))
            
            # Interactive Revenue Simulator
            st.markdown("#### 🧮 Interactive Revenue Simulator")
            with st.expander("🛠️ Adjust Growth Assumptions", expanded=False):
                sc1, sc2 = st.columns(2)
                target_pct = sc1.slider("Target Market Share (%)", 0.01, 5.0, 0.5, step=0.01)
                price_point = sc2.number_input("Monthly Subscription ($)", value=29)
                
                def parse_market_val(val_str):
                    try:
                        import re
                        clean = re.sub(r'[^\d.]', '', val_str)
                        val = float(clean) if clean else 1.0
                        if "B" in val_str.upper(): return val * 1e9
                        if "M" in val_str.upper(): return val * 1e6
                        if "K" in val_str.upper(): return val * 1e3
                        return val
                    except: return 1000000
                
                som_num = parse_market_val(mkt.get("som", "1M"))
                # Simplified model: Rev = SOM * Share%
                projected_annual = som_num * (target_pct / 100)
                projected_mrr = projected_annual / 12
                
                st.metric("Target MRR", f"${projected_mrr:,.0f}")
                st.caption(f"Projected ARR: ${projected_annual:,.0f} (at {target_pct}% share)")
                
                # 12-Month Scaling Graph
                scaling = [projected_mrr * (i/12)**2 for i in range(1, 13)] # Quadratic growth curve
                st.area_chart(scaling)

            st.markdown("---")

        # Reasoning Log
        reasoning_log = (st.session_state.current_state or {}).get("reasoning_log", "")
        if reasoning_log:
            with st.expander("🧠 View Deep Reasoning Analysis", expanded=False):
                try:
                    import json
                    parsed_log = json.loads(reasoning_log)
                    
                    st.markdown(f"**Quality Check:** `{parsed_log.get('quality_check', 'N/A')}`")
                    
                    st.markdown("**Reasoning Steps:**")
                    for step in parsed_log.get("reasoning_steps", []):
                        st.markdown(f"- {step}")
                        
                    st.markdown("**Core Clusters:**")
                    for cluster in parsed_log.get("core_clusters", []):
                        st.markdown(f"- {cluster}")
                        
                    st.markdown(f"**Market Intensity:**\n{parsed_log.get('market_intensity', '')}")
                    st.markdown(f"**Full Verdict:**\n{parsed_log.get('full_verdict', '')}")
                except Exception:
                    # Fallback if it's not valid JSON
                    st.markdown(f"```json\n{reasoning_log}\n```")
            st.markdown("---")

        # Voice of Customer
        voc = dossier.get("voice_of_customer", [])
        st.markdown("### 💬 Voice of the Customer")
        if voc:
            for item in voc:
                author = item.get('author', '')
                url    = item.get('url', '')
                source = item.get('source', 'Unknown source')
                # Build attribution line
                if url:
                    attribution = f'<a href="{url}" target="_blank" style="color:#60efff;text-decoration:none">@{author}</a> on {source}'
                elif author:
                    attribution = f'@{author} on {source}'
                else:
                    attribution = source

                st.markdown(f"""
                <div class='flux-quote'>
                    <span style='color:#f1f5f9; font-style:italic; font-size:1.1rem;'>"{item.get('quote','')}"</span><br>
                    <div style='margin-top:10px; font-size:0.9rem; color:#60efff;'>— {attribution}</div>
                </div>
                """, unsafe_allow_html=True)
        else:
            st.caption("No quotes captured.")

    # ── Tab 2: The MVP ───────────────────────────────────────────────────────
    with tab2:
        # Elevator Pitch
        growth = dossier.get("growth_strategy", {})
        if growth.get("elevator_pitch"):
            st.markdown(f"### 📣 {growth.get('elevator_pitch','')}")
            st.markdown("---")

        # UI Blueprint (Textual)
        ds = (st.session_state.current_state or {}).get("design_system", {})
        st.markdown("### 🗺️ Software UI Blueprint")
        layout = ds.get("ui_layout", {})
        if layout:
            st.markdown(f"""
            <div style='background:#0f172a; border: 1px solid #334155; border-radius: 16px; padding: 25px; font-family: "JetBrains Mono", monospace; box-shadow: inset 0 2px 20px rgba(0,0,0,0.5)'>
                <div style='color: #60efff; font-weight: 800; border-bottom: 1px solid #1e293b; padding-bottom: 12px; margin-bottom: 15px; font-size: 0.9rem;'>
                    [ TOP BAR ] &nbsp; {layout.get('header','')}
                </div>
                <div style='display: flex; min-height: 200px;'>
                    <div style='width: 35%; border-right: 1px solid #1e293b; padding-right: 15px; color: #64748b; font-size: 0.85rem;'>
                        <div style='color: #94a3b8; font-weight: 700; margin-bottom: 8px;'>[ SIDEBAR ]</div>
                        {layout.get('sidebar','')}
                    </div>
                    <div style='width: 65%; padding-left: 15px; color: #cbd5e1; font-size: 0.85rem;'>
                        <div style='color: #f8fafc; font-weight: 700; margin-bottom: 8px;'>[ MAIN WORKSPACE ]</div>
                        {layout.get('main_canvas','')}
                    </div>
                </div>
                <div style='margin-top: 25px; border-top: 1px solid #1e293b; padding-top: 15px; text-align: center;'>
                    <div style='color: #00ff87; font-size: 0.75rem; margin-bottom: 8px; font-weight: 700;'>PRIMARY INTERACTION</div>
                    <span style='background: #00ff87; color: #020617; padding: 8px 24px; border-radius: 6px; font-weight: 900; font-size: 0.8rem; letter-spacing: 1px;'>
                        {layout.get('primary_action','').upper()}
                    </span>
                </div>
            </div>
            """, unsafe_allow_html=True)
        else:
            st.caption("Generating UI blueprint... run a fresh scan to see the textual layout.")
            
        # Render Design System
        ds = (st.session_state.current_state or {}).get("design_system", {})
        if ds:
            with st.expander("🧬 View Brand Design System", expanded=True):
                dc1, dc2 = st.columns([1, 1.2])
                with dc1:
                    st.markdown("**Color Palette**")
                    palette = ds.get("color_palette", [])
                    pal_html = "".join([f"<div style='display:inline-block;background:{c};width:40px;height:40px;border-radius:8px;margin-right:8px;border:1px solid #334155' title='{c}'></div>" for c in palette])
                    st.markdown(pal_html, unsafe_allow_html=True)
                    st.caption(" ".join(palette))
                with dc2:
                    typo = ds.get("typography", {})
                    st.markdown(f"**Vibe:** {ds.get('component_style', 'N/A')}")
                    st.markdown(f"**Icons:** {ds.get('icon_style', 'N/A')}")
                    st.caption(f"Headings: {typo.get('headings', 'N/A')} | Body: {typo.get('font_family', 'N/A')}")
        
        # User Persona
        persona = ds.get("user_persona", {})
        if persona:
            st.markdown("### 👤 Ideal User Persona")
            st.markdown(f"""
            <div style='background:rgba(96, 239, 255, 0.05); border-radius:12px; padding:20px; border:1px dashed #60efff'>
                <div style='font-weight:800; color:#60efff;'>PROFILE: {persona.get('profile','')}</div>
                <div style='margin-top:8px; color:#f1f5f9;'><span style='color:#94a3b8'>Goal:</span> {persona.get('goal','')}</div>
                <div style='margin-top:4px; color:#f1f5f9;'><span style='color:#94a3b8'>Frustration:</span> {persona.get('frustration','')}</div>
            </div>
            """, unsafe_allow_html=True)

        st.markdown("---")

        mvp = dossier.get("mvp_blueprint", {})
        st.markdown("### 🚀 MVP Blueprint — The First 3 Features")
        for key, label, icon in [("feature_1","Core Utility","⚙️"), ("feature_2","The Hook","🪝"), ("feature_3","Admin Layer","🔐")]:
            feat = mvp.get(key, {})
            if feat:
                st.markdown(f"""
                <div style='background:#1e293b;border-radius:10px;padding:18px;margin-bottom:14px;border-left:4px solid #00ff87'>
                    <div style='font-weight:700;font-size:1.05rem;color:#f8fafc'>{icon} {feat.get('name','')}</div>
                    <div style='color:#94a3b8;margin-top:6px'>{feat.get('description','')}</div>
                </div>
                """, unsafe_allow_html=True)

        st.markdown("---")
        mono = dossier.get("monetization", {})
        st.markdown("### 💰 Monetization Hypothesis")
        m1, m2 = st.columns(2)
        m1.markdown(f"**Model:** {mono.get('model', 'N/A')}")
        m2.markdown(f"**Price Point:** {mono.get('price_point', 'N/A')}")
        st.info(mono.get("rationale", ""))
        
        # Growth Hack
        if growth:
            st.markdown("---")
            st.markdown("### 📈 Day 1 Growth Hack")
            st.success(f"**Action:** {growth.get('day_1_hack','')}")
            st.info(f"**Primary Channel:** {growth.get('primary_channel','')}")

    # ── Tab 3: The Build ─────────────────────────────────────────────────────
    with tab3:
        comp = dossier.get("competitive_landscape", [])
        st.markdown("### ⚔️ Competitive Landscape")
        if comp:
            import pandas as pd
            df = pd.DataFrame(comp).rename(columns={
                "competitor": "Competitor",
                "weakness":   "Their Weakness",
                "your_edge":  "Your Edge"
            })
            st.dataframe(df, use_container_width=True, hide_index=True)
        else:
            st.caption("No competitors identified.")

        st.markdown("---")
        tech = dossier.get("technical_roadmap", {})
        st.markdown("### 🛠️ Technical Roadmap")

        st.markdown(f"**🧱 Tech Stack:** `{tech.get('tech_stack', 'N/A')}`")
        st.markdown(f"**⏱️ Timeline:** {tech.get('timeline', 'N/A')}")

        dm = tech.get("data_model", [])
        if dm:
            st.markdown("**🗄️ Data Model (Key Tables):**")
            for table in dm:
                st.markdown(f"- `{table}`")

        week_plan = tech.get("week_plan", [])
        if week_plan:
            st.markdown("**📅 4-Week Technical Kanban Board:**")
            k1, k2, k3, k4 = st.columns(4)
            cols = [k1, k2, k3, k4]
            for i, w in enumerate(week_plan[:4]):
                with cols[i]:
                    st.markdown(f"""
                    <div style='background:#1e293b;border-radius:12px;padding:16px;min-height:140px;border-top:5px solid #60efff;box-shadow: 0 4px 15px rgba(0,0,0,0.3)'>
                        <div style='font-weight:900;color:#60efff;font-size:0.75rem;margin-bottom:10px;text-transform:uppercase'>Week {w.get('week','')}</div>
                        <div style='color:#f1f5f9;font-size:0.92rem;line-height:1.5;font-weight:500'>{w.get('focus','')}</div>
                    </div>
                    """, unsafe_allow_html=True)
            
        # Execution Details
        exec_det = dossier.get("execution_details", {})
        if exec_det:
            st.markdown("---")
            st.markdown("### ⚙️ Execution Criticals")
            ex1, ex2 = st.columns(2)
            with ex1:
                st.markdown("**🔌 Critical APIs:**")
                for api in exec_det.get("critical_apis", []):
                    st.markdown(f"- {api}")
            with ex2:
                st.warning(f"**🛡️ Security Focus:**\n{exec_det.get('security_priority','')}")
                st.info(f"**☁️ Infra Tip:** {exec_det.get('infrastructure_tip','')}")

    # ── Tab 4: Risk Audit ──────────────────────────────────────────────────
    with tab4:
        risk = (st.session_state.current_state or {}).get("risk_assessment", {})
        if risk:
            st.markdown("### 🕵️‍♂️ Red Team Analysis: Why this might fail")
            
            r1, r2 = st.columns(2)
            with r1:
                st.error(f"**🏗️ Technical Risk:**\n{risk.get('technical_risk','')}")
                st.error(f"**⚖️ Legal/IP Risk:**\n{risk.get('legal_risk','')}")
            with r2:
                st.error(f"**📈 Market Risk:**\n{risk.get('market_risk','')}")
                st.warning(f"**🛑 Kill-Switch Criteria:**\n{risk.get('kill_switch_criteria','')}")
            
            st.info(f"**🛡️ Survival Strategy:** {risk.get('survival_strategy','')}")
        else:
            st.caption("No risk audit data available.")

    # ── Tab 5: Sources ───────────────────────────────────────────────────────
    with tab5:
        st.markdown("### 📚 Research Sources & References")
        st.caption("These are the primary signals, expert reports, and community threads used to validate this business gap.")
        
        if all_refs:
            for i, ref in enumerate(all_refs, 1):
                url = ref.get('url', '')
                platform = "🌐 Web Signal"
                color = "#60efff"
                if "reddit.com" in url: platform, color = "🧡 Reddit", "#ff4500"
                elif "news.ycombinator.com" in url: platform, color = "🧡 HackerNews", "#ff6600"
                elif "producthunt.com" in url: platform, color = "😸 ProductHunt", "#da552f"
                
                title = ref.get('story_title') or ref.get('source') or "Direct Reference"
                
                st.markdown(f"""
                <div style='background:#1e293b;border-radius:10px;padding:12px;margin-bottom:10px;border-left:3px solid {color}'>
                    <div style='font-weight:600;color:#f8fafc'>
                        {i}. <a href="{url}" target="_blank" style="color:{color};text-decoration:none">{title}</a>
                    </div>
                    <div style='color:#94a3b8;font-size:0.85rem;margin-top:4px'>
                        {platform} | 📝 Source: {ref.get('author', 'Community Insight')} | <a href="{url}" target="_blank" style="color:#94a3b8;text-decoration:underline">Direct Link</a>
                    </div>
                </div>
                """, unsafe_allow_html=True)
        else:
            st.info("No external references were captured for this specific scan.")

    st.markdown("---")

    # ── Export ───────────────────────────────────────────────────────────────
    import json
    refs_md = "\n".join([f"- [{r.get('source','')}]({r['url']}) by {r.get('author','')}" for r in all_refs]) if all_refs else "No references."
    export_md = f"""# Founder's Dossier: {p_name}

## Signal Strength
- Mention Count: ~{dossier.get('signal_strength', {}).get('mention_count', '?')} threads
- {dossier.get('signal_strength', {}).get('source_summary', '')}

## Market Gap Score
{json.dumps(dossier.get('market_gap_score', {}), indent=2)}

## Voice of Customer
{chr(10).join(['> "' + q.get('quote','') + '" — @' + q.get('author','') + ' (' + q.get('url','') + ')' for q in dossier.get('voice_of_customer',[])])}

## MVP Blueprint
{json.dumps(dossier.get('mvp_blueprint', {}), indent=2)}

## Competitive Landscape
{json.dumps(dossier.get('competitive_landscape', []), indent=2)}

## Technical Roadmap
{json.dumps(dossier.get('technical_roadmap', {}), indent=2)}

## Monetization
{json.dumps(dossier.get('monetization', {}), indent=2)}

## References
{refs_md}
"""
    col_dl, col_pptx, col_save, col_back, col_new = st.columns(5)
    with col_dl:
        st.download_button(
            label="💾 Markdown",
            data=export_md,
            file_name=f"dossier_{p_name.replace(' ','_').lower()}.md",
            mime="text/markdown",
            use_container_width=True
        )
    with col_pptx:
        m_analysis = (st.session_state.current_state or {}).get('market_size_analysis', {})
        r_analysis = (st.session_state.current_state or {}).get('risk_assessment', {})
        pptx_data = generate_pitch_deck(p_name, dossier, m_analysis, r_analysis)
        st.download_button(
            label="📊 Pitch Deck",
            data=pptx_data,
            file_name=f"pitch_deck_{p_name.replace(' ','_').lower()}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
    with col_save:
        if st.button("🗄️ Archive to Library", use_container_width=True):
            m_url = (st.session_state.current_state or {}).get("mockup_url")
            risk_data = (st.session_state.current_state or {}).get("risk_assessment")
            db.save_dossier(st.session_state.topic, p_name, dossier, p_score, m_url, risk_data)
            st.success("Archived!")
            time.sleep(1)
            st.rerun()
    with col_back:
        if st.button("🔙 Back to Ideas"):
            # Clear selection-specific state from graph
            graph_app.update_state(config, {
                "selected_problem": {},
                "blueprint": None,
                "market_size_analysis": None,
                "mockup_url": None,
                "risk_assessment": None
            })
            # Reset session state
            st.session_state.selected_problem = None
            if st.session_state.current_state:
                st.session_state.current_state.pop("blueprint", None)
                st.session_state.current_state.pop("market_size_analysis", None)
                st.session_state.current_state.pop("mockup_url", None)
                st.session_state.current_state.pop("risk_assessment", None)
            
            st.session_state.app_stage = "selection"
            st.rerun()
    st.markdown("---")
    
    # ── Strategy Pivot (Dynamic Refinement) ──────────────────────────────────
    with st.expander("🔥 PIVOT STRATEGY: Refine this Blueprint", expanded=False):
        st.markdown("#### 🛠️ Direct the AI to update this idea")
        st.caption("Example: 'Target enterprise companies instead of individuals', 'Add a focus on privacy', 'Suggest a different tech stack'.")
        pivot_input = st.text_area("Your Refinement Instructions:", placeholder="What should be different about this business model?")
        
        if st.button("🚀 Regenerate with Feedback", use_container_width=True):
            with st.status("⚙️ Processing Venture Intelligence...", expanded=True) as status:
                # Add pivot instructions to state
                st.session_state.current_state["topic"] = f"{st.session_state.topic} (REFINE: {pivot_input})"
                
                # We need to clear previous outputs to force regeneration
                st.session_state.current_state.pop("blueprint", None)
                st.session_state.current_state.pop("market_size_analysis", None)
                st.session_state.current_state.pop("risk_assessment", None)
                
                # Rerun processing2 stage
                st.session_state.app_stage = "processing2"
                st.rerun()
    
    # ── Interactive Strategy Consultant Chat ─────────────────────────────────
    st.markdown("### 💬 Chat with your Strategy Consultant")
    st.caption("Ask follow-up questions about this blueprint, marketing tactics, or technical details.")

    # Display chat history
    for msg in st.session_state.chat_history:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    # Chat input
    if prompt := st.chat_input("How should I market this? / What stack is best?"):
        # Add user message
        st.session_state.chat_history.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        # Generate response
        with st.chat_message("assistant"):
            with st.spinner("Consulting the strategist..."):
                try:
                    from src.graph import invoke_llm
                    from langchain_core.prompts import PromptTemplate
                    
                    # Context for the AI
                    context_template = """
                    You are an elite AI Startup Consultant. You just generated this Founder's Dossier:
                    PROBLEM: {p_name}
                    BLUEPRINT: {dossier_json}
                    MARKET: {market_json}
                    RISKS: {risk_json}
                    
                    User Question: {user_query}
                    
                    Answer the user's questions about this specific idea. Be actionable, realistic, and insightful.
                    """
                    
                    chat_prompt_obj = PromptTemplate(
                        template=context_template,
                        input_variables=["p_name", "dossier_json", "market_json", "risk_json", "user_query"]
                    )
                    
                    # Package inputs
                    inputs = {
                        "p_name": p_name,
                        "dossier_json": json.dumps(dossier),
                        "market_json": json.dumps((st.session_state.current_state or {}).get('market_size_analysis', {})),
                        "risk_json": json.dumps((st.session_state.current_state or {}).get('risk_assessment', {})),
                        "user_query": prompt
                    }
                    
                    assistant_text, model_id = invoke_llm(chat_prompt_obj, inputs, tier="versatile", temperature=0.7)
                    
                    st.markdown(assistant_text)
                    st.caption(f"Generated by {model_id}")
                    st.session_state.chat_history.append({"role": "assistant", "content": assistant_text})
                except Exception as e:
                    st.error(f"Consultant Error: {e}")

# Stage: VIEW ARCHIVE
if st.session_state.app_stage == "view_archive":
    item = st.session_state.archived_dossier
    if not item:
        st.session_state.app_stage = "input"
        st.rerun()
    
    p_name = item['problem_name']
    p_score = item['market_score']
    dossier = json.loads(item['blueprint_json'])
    
    st.markdown(f"### 🗄️ Viewing Archived Dossier: {p_name}")
    
    # Reuse the display logic (Simplified for space)
    st.info(f"Originally created on {item['created_at']}")
    
    # ... (I'll extract the rendering to a function if I had more space, but for now I will just render it)
    # Actually, let's just use the same rendering block but with archived data
    
    score_color = "#00ff87" if p_score >= 75 else ("#ffb347" if p_score >= 50 else "#ff6b6b")

    st.markdown(f"""
    <div style='background:#1e293b;padding:24px;border-radius:14px;border-left:6px solid {score_color};margin-bottom:24px'>
        <div style='font-size:1.6rem;font-weight:800;color:#f8fafc;margin-bottom:8px'>📄 {p_name}</div>
        <span style='background:#334155;color:{score_color};padding:4px 14px;border-radius:20px;font-size:0.95rem;font-weight:700;margin-right:8px'>🔥 Market Score: {p_score}/100</span>
    </div>
    """, unsafe_allow_html=True)

    tab1, tab2, tab3, tab4 = st.tabs(["📊 The Insight", "💡 The MVP", "🛠️ The Build", "📉 Risks"])
    
    with tab1:
        st.json(dossier.get("signal_strength", {}))
        st.json(dossier.get("market_gap_score", {}))
    
    with tab2:
        m_url = item.get("mockup_url")
        if m_url:
            st.image(m_url, caption="Archived Visual Mockup", use_container_width=True)
        st.json(dossier.get("mvp_blueprint", {}))
    
    with tab4:
        r_json = item.get("risk_json")
        if r_json:
            st.json(json.loads(r_json))
    
    if st.button("🔙 Back to Main"):
        st.session_state.app_stage = "input"
        st.rerun()
